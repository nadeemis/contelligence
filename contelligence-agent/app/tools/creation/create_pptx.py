"""PowerPoint presentation creation and editing tool.

Creates professional, beautifully styled PPTX presentations from a
declarative slide specification.  Supports built-in themes, custom
colours, template files, and fine-grained element placement.
"""

import base64
import io
import logging
import pathlib
from typing import Any, Literal, Self

from pydantic import BaseModel, Field, model_validator

from app.core.tool_registry import define_tool

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_SLIDE_WIDTH_IN = 13.333
_SLIDE_HEIGHT_IN = 7.5

# Margins / padding (inches)
_MARGIN_LEFT = 0.8
_MARGIN_RIGHT = 0.8
_MARGIN_TOP = 0.6
_MARGIN_BOTTOM = 0.5
_CONTENT_TOP = 1.7  # below title area
_CONTENT_WIDTH = _SLIDE_WIDTH_IN - _MARGIN_LEFT - _MARGIN_RIGHT

# ---------------------------------------------------------------------------
# Theme palettes
# ---------------------------------------------------------------------------

_THEMES: dict[str, dict[str, Any]] = {
    "corporate": {
        "bg": "FFFFFF",
        "title": "1E3A5F",
        "text": "333333",
        "accent1": "0078D4",
        "accent2": "107C10",
        "accent3": "D83B01",
        "subtitle": "5A6A7A",
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": 36,
        "subtitle_size": 20,
        "body_size": 18,
        "bullet_size": 17,
        "heading_size": 22,
    },
    "modern_dark": {
        "bg": "1A1A2E",
        "title": "FFFFFF",
        "text": "E0E0E0",
        "accent1": "00D2FF",
        "accent2": "7C3AED",
        "accent3": "10B981",
        "subtitle": "A0AEC0",
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_size": 36,
        "subtitle_size": 20,
        "body_size": 18,
        "bullet_size": 17,
        "heading_size": 22,
    },
    "creative": {
        "bg": "F8F9FA",
        "title": "2D2D2D",
        "text": "4A4A4A",
        "accent1": "FF6B6B",
        "accent2": "48BFE3",
        "accent3": "F7D794",
        "subtitle": "6C757D",
        "title_font": "Segoe UI Semibold",
        "body_font": "Segoe UI",
        "title_size": 38,
        "subtitle_size": 20,
        "body_size": 18,
        "bullet_size": 17,
        "heading_size": 22,
    },
    "minimal": {
        "bg": "FFFFFF",
        "title": "111111",
        "text": "555555",
        "accent1": "000000",
        "accent2": "888888",
        "accent3": "DDDDDD",
        "subtitle": "777777",
        "title_font": "Calibri Light",
        "body_font": "Calibri",
        "title_size": 40,
        "subtitle_size": 18,
        "body_size": 16,
        "bullet_size": 15,
        "heading_size": 20,
    },
}

# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------


class ThemeOverrides(BaseModel):
    """Optional per-field overrides applied on top of a named theme."""

    bg_color: str | None = Field(None, description="Background colour as 6-digit hex, e.g. 'FF0000'.")
    title_color: str | None = Field(None, description="Title text colour (hex).")
    text_color: str | None = Field(None, description="Body text colour (hex).")
    accent_color: str | None = Field(None, description="Primary accent colour (hex).")
    title_font: str | None = Field(None, description="Font family for titles.")
    body_font: str | None = Field(None, description="Font family for body text.")


class ElementSpec(BaseModel):
    """A custom element placed at an absolute position on a slide."""

    type: Literal["text", "image", "shape"] = Field(
        ..., description="Element kind: 'text' for a text box, 'image' for a picture, 'shape' for a geometric shape."
    )
    left: float = Field(..., description="Left edge position in inches from the slide left.")
    top: float = Field(..., description="Top edge position in inches from the slide top.")
    width: float = Field(..., description="Element width in inches.")
    height: float = Field(..., description="Element height in inches.")

    # ---- text properties ----
    text: str | None = Field(None, description="Text content (for type='text').")
    font_size: float | None = Field(None, description="Font size in points.")
    font_color: str | None = Field(None, description="Font colour as hex, e.g. 'FFFFFF'.")
    font_bold: bool = Field(False, description="Bold text.")
    font_italic: bool = Field(False, description="Italic text.")
    font_name: str | None = Field(None, description="Font family name.")
    alignment: Literal["left", "center", "right"] | None = Field(None, description="Horizontal text alignment.")
    vertical_alignment: Literal["top", "middle", "bottom"] | None = Field(None, description="Vertical text alignment within the box.")
    line_spacing: float | None = Field(None, description="Line spacing multiplier (e.g. 1.2).")

    # ---- image properties ----
    image_b64: str | None = Field(None, description="Base64-encoded image data (PNG/JPG).")

    # ---- shape properties ----
    shape_type: Literal[
        "rectangle", "rounded_rectangle", "circle", "oval",
        "triangle", "diamond", "pentagon", "hexagon",
        "arrow_right", "arrow_left", "arrow_up", "arrow_down",
        "chevron_right", "star_5", "star_6", "heart",
        "line_horizontal", "line_vertical",
    ] | None = Field(None, description="Auto-shape type (for type='shape').")
    fill_color: str | None = Field(None, description="Shape fill colour (hex). Use 'transparent' for no fill.")
    border_color: str | None = Field(None, description="Shape border colour (hex).")
    border_width: float | None = Field(None, description="Shape border width in points.")
    text_inside: str | None = Field(None, description="Text rendered inside the shape.")
    text_color: str | None = Field(None, description="Colour of text inside the shape (hex).")
    text_size: float | None = Field(None, description="Font size for text inside the shape (points).")
    text_bold: bool = Field(False, description="Bold text inside the shape.")
    shadow: bool = Field(False, description="Add a drop shadow to the shape.")
    rotation: float | None = Field(None, description="Rotation angle in degrees (clockwise).")


class ContentBlock(BaseModel):
    """A block of content used in the ``content_blocks`` shorthand.

    LLMs often prefer producing a flat ``content_blocks`` list rather than
    mapping content to ``body``/``left_column``/``right_column`` individually.
    This model captures that pattern and a model validator on :class:`SlideSpec`
    normalises it into the canonical fields the layout builders consume.
    """

    type: Literal["text", "image"] = Field("text", description="Block type.")
    value: str | None = Field(None, description="Text content for the block.")
    label: str | None = Field(None, description="Optional header/label for the block (used as column header in two-column layouts).")
    image_b64: str | None = Field(None, description="Base64-encoded image data (for type='image').")
    url: str | None = Field(None, description="URL to an image to download (for type='image').")


class SlideSpec(BaseModel):
    """Specification for a single slide."""

    layout: Literal[
        "title", "section", "content", "two_column",
        "image_left", "image_right", "image+text", "image_text",
        "comparison", "blank",
    ] = Field("content", description=(
        "Slide layout: 'title' for title/subtitle opener, 'section' for a divider, "
        "'content' for title + bullet body, 'two_column' for side-by-side columns, "
        "'image_left'/'image_right'/'image+text' for image + text, 'comparison' for vs-style, "
        "'blank' for a fully custom slide."
    ))

    title: str | None = Field(None, description="Slide title text.")
    subtitle: str | None = Field(None, description="Subtitle (mainly for 'title' and 'section' layouts).")
    body: str | list[str] | None = Field(None, description="Body text. A string for a paragraph, a list for bullet points.")
    bullets: list[str] | None = Field(None, description="Explicit bullet-point list (overrides body if both given).")

    # Two-column / comparison
    left_column: str | list[str] | None = Field(None, description="Left column content (string or bullet list).")
    right_column: str | list[str] | None = Field(None, description="Right column content (string or bullet list).")
    left_header: str | None = Field(None, description="Header for the left column (comparison/two_column layouts).")
    right_header: str | None = Field(None, description="Header for the right column.")

    # Image
    image_b64: str | None = Field(None, description="Base64-encoded image for image layouts (PNG/JPG).")
    image_url: str | None = Field(None, description="URL to an image to download for image layouts.")
    image_description: str | None = Field(None, description="Alt text / caption for the image.")

    # Speaker notes
    notes: str | None = Field(None, description="Speaker notes for this slide.")

    # Custom positioned elements (additive — placed on top of layout content)
    elements: list[ElementSpec] | None = Field(None, description="Custom positioned elements added on top of the layout.")

    # Per-slide overrides
    bg_color: str | None = Field(None, description="Override background colour for this slide (hex).")

    # Flexible content input (LLM-friendly alternative to body/bullets/left_column/right_column)
    content_blocks: list[ContentBlock] | None = Field(
        None,
        description=(
            "Flexible content blocks — an alternative to body/bullets/left_column/right_column. "
            "For 'content' layouts: a single block's 'value' becomes the body text. "
            "For 'two_column'/'comparison' layouts: two blocks map to left/right columns, "
            "with each block's 'label' becoming the column header."
        ),
    )

    @model_validator(mode="after")
    def _normalise_content_blocks(self) -> Self:
        """Map ``content_blocks`` into the canonical fields when they are empty.

        Also normalises layout aliases (``image+text`` → ``image_left``).
        """
        # Normalise layout aliases
        _LAYOUT_ALIASES: dict[str, str] = {
            "image+text": "image_left",
            "image_text": "image_left",
        }
        if self.layout in _LAYOUT_ALIASES:
            self.layout = _LAYOUT_ALIASES[self.layout]  # type: ignore[assignment]

        blocks = self.content_blocks
        if not blocks:
            return self

        # Extract image info from content_blocks (URL or base64)
        img_blocks = [b for b in blocks if b.type == "image"]
        if img_blocks:
            blk = img_blocks[0]
            if blk.url and self.image_url is None and self.image_b64 is None:
                self.image_url = blk.url
            elif blk.image_b64 and self.image_b64 is None:
                self.image_b64 = blk.image_b64

        if self.layout in ("content", "section", "image_left", "image_right"):
            # Single-body layouts: use the first text block as body
            if self.body is None and self.bullets is None:
                text_blocks = [b for b in blocks if b.type == "text" and b.value]
                if text_blocks:
                    self.body = text_blocks[0].value

        elif self.layout in ("two_column", "comparison"):
            # Two-column layouts: first block → left, second → right
            text_blocks = [b for b in blocks if b.type == "text"]
            if len(text_blocks) >= 2:
                if self.left_column is None:
                    self.left_column = text_blocks[0].value
                if self.left_header is None and text_blocks[0].label:
                    self.left_header = text_blocks[0].label
                if self.right_column is None:
                    self.right_column = text_blocks[1].value
                if self.right_header is None and text_blocks[1].label:
                    self.right_header = text_blocks[1].label
            elif len(text_blocks) == 1:
                if self.left_column is None:
                    self.left_column = text_blocks[0].value
                if self.left_header is None and text_blocks[0].label:
                    self.left_header = text_blocks[0].label

        elif self.layout == "title":
            # Title layout: first block's value can fill subtitle
            text_blocks = [b for b in blocks if b.type == "text" and b.value]
            if text_blocks and self.subtitle is None:
                self.subtitle = text_blocks[0].value

        return self


class CreatePptxParams(BaseModel):
    """Parameters for the create_pptx tool."""

    slides: list[SlideSpec] = Field(
        ...,
        description="Ordered list of slide specifications. Each slide is rendered according to its layout.",
        min_length=1,
    )

    theme: Literal["corporate", "modern_dark", "creative", "minimal"] = Field(
        "corporate",
        description="Named colour theme for the presentation.",
    )
    custom_theme: ThemeOverrides | None = Field(
        None,
        description="Optional overrides applied on top of the chosen theme (e.g. custom accent colour).",
    )

    # Template (optional — mutually exclusive sources)
    template_bytes_b64: str | None = Field(
        None, description="Base64-encoded PPTX template file to use as the starting point.",
    )
    template_local_path: str | None = Field(
        None, description="Local file path to a PPTX template.",
    )
    template_storage_account: str | None = Field(
        None, description="Azure Storage account for the template blob.",
    )
    template_container: str | None = Field(
        None, description="Azure Blob container for the template.",
    )
    template_path: str | None = Field(
        None, description="Blob path to the PPTX template.",
    )

    # Output destinations (all optional — base64 is always returned)
    output_local_path: str | None = Field(
        None, description="Save the generated PPTX to this local path.",
    )
    output_storage_account: str | None = Field(
        None, description="Azure Storage account for uploading the output.",
    )
    output_container: str | None = Field(
        None, description="Azure Blob container for the output.",
    )
    output_path: str | None = Field(
        None, description="Blob path for the output PPTX.",
    )
    output_filename: str | None = Field(
        "presentation.pptx",
        description="Descriptive filename for the output (used in response metadata).",
    )

    slide_width: float = Field(
        _SLIDE_WIDTH_IN,
        description="Slide width in inches (default 13.333 for 16:9 widescreen).",
    )
    slide_height: float = Field(
        _SLIDE_HEIGHT_IN,
        description="Slide height in inches (default 7.5 for 16:9 widescreen).",
    )


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _hex(h: str) -> "RGBColor":
    """Convert a 6-digit hex string to an ``RGBColor``."""
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]

    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _inches(val: float) -> int:
    """Convert inches to EMU."""
    from pptx.util import Inches  # type: ignore[import-untyped]

    return Inches(val)


def _pt(val: float) -> int:
    """Convert points to EMU."""
    from pptx.util import Pt  # type: ignore[import-untyped]

    return Pt(val)


def _resolve_theme(params: CreatePptxParams) -> dict[str, Any]:
    """Merge base theme with optional overrides."""
    base = dict(_THEMES[params.theme])
    if params.custom_theme:
        ov = params.custom_theme
        if ov.bg_color:
            base["bg"] = ov.bg_color
        if ov.title_color:
            base["title"] = ov.title_color
        if ov.text_color:
            base["text"] = ov.text_color
        if ov.accent_color:
            base["accent1"] = ov.accent_color
        if ov.title_font:
            base["title_font"] = ov.title_font
        if ov.body_font:
            base["body_font"] = ov.body_font
    return base


# ---------------------------------------------------------------------------
# Shape-type mapping
# ---------------------------------------------------------------------------

_SHAPE_MAP: dict[str, int] = {}


def _get_shape_map() -> dict[str, int]:
    if _SHAPE_MAP:
        return _SHAPE_MAP
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]

    mapping = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "circle": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "arrow_right": MSO_SHAPE.RIGHT_ARROW,
        "arrow_left": MSO_SHAPE.LEFT_ARROW,
        "arrow_up": MSO_SHAPE.UP_ARROW,
        "arrow_down": MSO_SHAPE.DOWN_ARROW,
        "chevron_right": MSO_SHAPE.CHEVRON,
        "star_5": MSO_SHAPE.STAR_5_POINT,
        "star_6": MSO_SHAPE.STAR_6_POINT,
        "heart": MSO_SHAPE.HEART,
        "line_horizontal": MSO_SHAPE.RECTANGLE,
        "line_vertical": MSO_SHAPE.RECTANGLE,
    }
    _SHAPE_MAP.update(mapping)
    return _SHAPE_MAP


# ---------------------------------------------------------------------------
# Slide background
# ---------------------------------------------------------------------------


def _set_slide_bg(slide: Any, colour_hex: str) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex(colour_hex)


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


def _add_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str = "Calibri",
    font_size: float = 18,
    font_color: str = "333333",
    bold: bool = False,
    italic: bool = False,
    alignment: str = "left",
    vertical: str = "top",
    line_spacing: float | None = None,
) -> Any:
    """Add a styled text box and return the shape."""
    from pptx.util import Emu  # type: ignore[import-untyped]
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import-untyped]

    txBox = slide.shapes.add_textbox(
        _inches(left), _inches(top), _inches(width), _inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    _valign_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.paragraphs[0].alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    }.get(alignment, PP_ALIGN.LEFT)

    try:
        txBox.vertical_anchor = _valign_map.get(vertical, MSO_ANCHOR.TOP)  # type: ignore[assignment]
    except Exception:
        pass

    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = _pt(font_size)
    run.font.color.rgb = _hex(font_color)
    run.font.bold = bold
    run.font.italic = italic

    if line_spacing is not None:
        from pptx.util import Pt as PtUnit
        p.line_spacing = PtUnit(font_size * line_spacing)

    return txBox


def _add_bullets(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    font_name: str = "Calibri",
    font_size: float = 17,
    font_color: str = "333333",
    bullet_color: str | None = None,
    line_spacing: float = 1.4,
) -> Any:
    """Add a text box with bullet points."""
    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    txBox = slide.shapes.add_textbox(
        _inches(left), _inches(top), _inches(width), _inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.space_after = _pt(4)

        # Set line spacing
        from pptx.util import Pt as PtUnit
        p.line_spacing = PtUnit(font_size * line_spacing)

        run = p.runs[0] if p.runs else p.add_run()
        run.text = item
        run.font.name = font_name
        run.font.size = _pt(font_size)
        run.font.color.rgb = _hex(font_color)

        # Enable bullet
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        # Remove existing bullets first
        for existing in pPr.findall(qn("a:buChar")):
            pPr.remove(existing)
        for existing in pPr.findall(qn("a:buNone")):
            pPr.remove(existing)
        pPr.append(buChar)

        # Bullet colour
        if bullet_color:
            buClr = pPr.makeelement(qn("a:buClr"), {})
            srgb = buClr.makeelement(qn("a:srgbClr"), {"val": bullet_color})
            buClr.append(srgb)
            for existing in pPr.findall(qn("a:buClr")):
                pPr.remove(existing)
            pPr.append(buClr)

        # Bullet size (relative to text)
        buSzPct = pPr.makeelement(qn("a:buSzPct"), {"val": "100000"})
        for existing in pPr.findall(qn("a:buSzPct")):
            pPr.remove(existing)
        pPr.append(buSzPct)

        # Indent
        pPr.set("marL", str(int(228600)))  # 0.25 inch
        pPr.set("indent", str(int(-228600)))

    return txBox


# ---------------------------------------------------------------------------
# Decoration helpers
# ---------------------------------------------------------------------------


def _add_accent_bar(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    colour: str,
) -> Any:
    """Add a thin coloured rectangle (accent bar)."""
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _inches(left), _inches(top), _inches(width), _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(colour)
    shape.line.fill.background()
    return shape


def _add_shape_element(
    slide: Any,
    spec: ElementSpec,
    theme: dict[str, Any],
) -> Any:
    """Add a geometric shape from an ElementSpec."""
    shape_map = _get_shape_map()
    shape_kind = spec.shape_type or "rectangle"
    mso_type = shape_map.get(shape_kind, shape_map["rectangle"])

    # Thin lines are just very thin rectangles
    h = spec.height
    w = spec.width
    if shape_kind == "line_horizontal":
        h = max(h, 0.02)
    elif shape_kind == "line_vertical":
        w = max(w, 0.02)

    shape = slide.shapes.add_shape(
        mso_type,
        _inches(spec.left), _inches(spec.top), _inches(w), _inches(h),
    )

    # Fill
    if spec.fill_color and spec.fill_color.lower() == "transparent":
        shape.fill.background()
    elif spec.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(spec.fill_color)
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(theme["accent1"])

    # Border
    if spec.border_color:
        shape.line.color.rgb = _hex(spec.border_color)
        shape.line.width = _pt(spec.border_width or 1)
    else:
        shape.line.fill.background()

    # Text inside shape
    if spec.text_inside:
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import-untyped]

        tf = shape.text_frame
        tf.word_wrap = True
        try:
            shape.text_frame.auto_size = None  # type: ignore[assignment]
            shape.vertical_anchor = MSO_ANCHOR.MIDDLE  # type: ignore[assignment]
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.text = spec.text_inside
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.text = spec.text_inside
        run.font.name = spec.font_name or theme["body_font"]
        run.font.size = _pt(spec.text_size or 14)
        run.font.color.rgb = _hex(spec.text_color or theme["title"])
        run.font.bold = spec.text_bold

    # Rotation
    if spec.rotation:
        shape.rotation = spec.rotation

    # Shadow (via XML manipulation)
    if spec.shadow:
        _apply_shadow(shape)

    return shape


def _apply_shadow(shape: Any) -> None:
    """Apply a subtle outer shadow to a shape via OpenXML."""
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "50800",
        "dist": "38100",
        "dir": "5400000",
        "rotWithShape": "0",
    })
    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "40000"})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    # Remove existing effects
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)
    spPr.append(effectLst)


def _add_image_from_b64(
    slide: Any,
    b64_data: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Any:
    """Add an image from base64 data."""
    img_bytes = base64.b64decode(b64_data)
    stream = io.BytesIO(img_bytes)
    return slide.shapes.add_picture(
        stream, _inches(left), _inches(top), _inches(width), _inches(height),
    )


# ---------------------------------------------------------------------------
# Custom element placement
# ---------------------------------------------------------------------------


def _place_element(slide: Any, spec: ElementSpec, theme: dict[str, Any]) -> None:
    """Place a single custom element on a slide."""
    if spec.type == "text":
        _add_textbox(
            slide,
            spec.left, spec.top, spec.width, spec.height,
            spec.text or "",
            font_name=spec.font_name or theme["body_font"],
            font_size=spec.font_size or theme["body_size"],
            font_color=spec.font_color or theme["text"],
            bold=spec.font_bold,
            italic=spec.font_italic,
            alignment=spec.alignment or "left",
            vertical=spec.vertical_alignment or "top",
            line_spacing=spec.line_spacing,
        )
    elif spec.type == "image" and spec.image_b64:
        _add_image_from_b64(
            slide, spec.image_b64,
            spec.left, spec.top, spec.width, spec.height,
        )
    elif spec.type == "shape":
        _add_shape_element(slide, spec, theme)


# ---------------------------------------------------------------------------
# Layout builders
# ---------------------------------------------------------------------------


def _build_title_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a title slide with centred title, subtitle, and accent bar."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    sw = prs.slide_width.inches if hasattr(prs.slide_width, "inches") else _SLIDE_WIDTH_IN

    # Accent bar — centred, above title
    bar_w = 3.0
    _add_accent_bar(
        slide,
        left=(sw - bar_w) / 2, top=2.2, width=bar_w, height=0.06,
        colour=theme["accent1"],
    )

    # Title
    if spec.title:
        _add_textbox(
            slide, left=1.0, top=2.5, width=sw - 2.0, height=1.4,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"] + 4,
            font_color=theme["title"],
            bold=True,
            alignment="center",
            vertical="middle",
        )

    # Subtitle
    if spec.subtitle:
        _add_textbox(
            slide, left=2.0, top=4.1, width=sw - 4.0, height=0.9,
            text=spec.subtitle,
            font_name=theme["body_font"],
            font_size=theme["subtitle_size"],
            font_color=theme["subtitle"],
            alignment="center",
            vertical="top",
        )

    # Bottom accent line
    _add_accent_bar(
        slide,
        left=(sw - 5.0) / 2, top=5.3, width=5.0, height=0.03,
        colour=theme["accent1"],
    )

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _build_section_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a section divider slide with a bold vertical accent strip."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    # Vertical accent strip on the left
    _add_accent_bar(slide, left=0, top=0, width=0.35, height=7.5, colour=theme["accent1"])

    # Section title
    if spec.title:
        _add_textbox(
            slide, left=1.2, top=2.4, width=10.0, height=1.5,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"] + 2,
            font_color=theme["title"],
            bold=True,
            alignment="left",
            vertical="bottom",
        )

    # Accent underline
    _add_accent_bar(slide, left=1.2, top=4.0, width=4.0, height=0.05, colour=theme["accent1"])

    # Subtitle
    if spec.subtitle:
        _add_textbox(
            slide, left=1.2, top=4.3, width=10.0, height=0.8,
            text=spec.subtitle,
            font_name=theme["body_font"],
            font_size=theme["subtitle_size"],
            font_color=theme["subtitle"],
            alignment="left",
        )

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _build_content_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a standard content slide: title + body / bullets."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    # Title
    if spec.title:
        _add_textbox(
            slide, left=_MARGIN_LEFT, top=_MARGIN_TOP, width=_CONTENT_WIDTH, height=0.7,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"],
            font_color=theme["title"],
            bold=True,
        )
        # Accent underline
        _add_accent_bar(
            slide, left=_MARGIN_LEFT, top=1.35, width=2.0, height=0.04,
            colour=theme["accent1"],
        )

    # Body content
    items = spec.bullets or (spec.body if isinstance(spec.body, list) else None)
    body_top = _CONTENT_TOP
    body_height = _SLIDE_HEIGHT_IN - body_top - _MARGIN_BOTTOM

    if items:
        _add_bullets(
            slide, _MARGIN_LEFT, body_top, _CONTENT_WIDTH, body_height,
            items,
            font_name=theme["body_font"],
            font_size=theme["bullet_size"],
            font_color=theme["text"],
            bullet_color=theme["accent1"],
        )
    elif spec.body and isinstance(spec.body, str):
        _add_textbox(
            slide, _MARGIN_LEFT, body_top, _CONTENT_WIDTH, body_height,
            spec.body,
            font_name=theme["body_font"],
            font_size=theme["body_size"],
            font_color=theme["text"],
            line_spacing=1.4,
        )

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _build_two_column_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a two-column layout with optional column headers."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    # Title
    if spec.title:
        _add_textbox(
            slide, left=_MARGIN_LEFT, top=_MARGIN_TOP, width=_CONTENT_WIDTH, height=0.7,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"],
            font_color=theme["title"],
            bold=True,
        )
        _add_accent_bar(
            slide, left=_MARGIN_LEFT, top=1.35, width=2.0, height=0.04,
            colour=theme["accent1"],
        )

    col_width = (_CONTENT_WIDTH - 0.5) / 2  # gap between columns
    left_x = _MARGIN_LEFT
    right_x = _MARGIN_LEFT + col_width + 0.5
    header_top = _CONTENT_TOP
    content_top = _CONTENT_TOP + 0.55
    col_height = _SLIDE_HEIGHT_IN - content_top - _MARGIN_BOTTOM

    # Subtle vertical divider
    divider_x = _MARGIN_LEFT + col_width + 0.22
    _add_accent_bar(slide, left=divider_x, top=_CONTENT_TOP, width=0.02, height=col_height + 0.4, colour=theme["accent3"])

    # Left header
    if spec.left_header:
        _add_textbox(
            slide, left_x, header_top, col_width, 0.45,
            spec.left_header,
            font_name=theme["title_font"],
            font_size=theme["heading_size"],
            font_color=theme["accent1"],
            bold=True,
        )

    # Right header
    if spec.right_header:
        _add_textbox(
            slide, right_x, header_top, col_width, 0.45,
            spec.right_header,
            font_name=theme["title_font"],
            font_size=theme["heading_size"],
            font_color=theme["accent1"],
            bold=True,
        )

    # Left content
    _render_column(slide, spec.left_column, left_x, content_top, col_width, col_height, theme)

    # Right content
    _render_column(slide, spec.right_column, right_x, content_top, col_width, col_height, theme)

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _render_column(
    slide: Any,
    content: str | list[str] | None,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, Any],
) -> None:
    """Render column content as bullets or paragraph."""
    if content is None:
        return
    if isinstance(content, list):
        _add_bullets(
            slide, left, top, width, height,
            content,
            font_name=theme["body_font"],
            font_size=theme["bullet_size"],
            font_color=theme["text"],
            bullet_color=theme["accent1"],
        )
    else:
        _add_textbox(
            slide, left, top, width, height,
            content,
            font_name=theme["body_font"],
            font_size=theme["body_size"],
            font_color=theme["text"],
            line_spacing=1.4,
        )


def _build_image_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
    image_side: Literal["left", "right"],
) -> None:
    """Build an image + text slide with the image on the specified side."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    # Title (full width at top)
    if spec.title:
        _add_textbox(
            slide, left=_MARGIN_LEFT, top=_MARGIN_TOP, width=_CONTENT_WIDTH, height=0.7,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"],
            font_color=theme["title"],
            bold=True,
        )
        _add_accent_bar(
            slide, left=_MARGIN_LEFT, top=1.35, width=2.0, height=0.04,
            colour=theme["accent1"],
        )

    img_width = 5.2
    img_height = 4.2
    text_width = _CONTENT_WIDTH - img_width - 0.6
    img_top = _CONTENT_TOP + 0.2
    text_top = _CONTENT_TOP + 0.2

    if image_side == "left":
        img_left = _MARGIN_LEFT
        text_left = _MARGIN_LEFT + img_width + 0.6
    else:
        text_left = _MARGIN_LEFT
        img_left = _MARGIN_LEFT + text_width + 0.6

    # Image placeholder or actual image
    if spec.image_b64:
        _add_image_from_b64(slide, spec.image_b64, img_left, img_top, img_width, img_height)
    else:
        # Placeholder rectangle with icon-like text
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import-untyped]

        ph = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            _inches(img_left), _inches(img_top), _inches(img_width), _inches(img_height),
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = _hex(theme["accent3"])
        ph.line.fill.background()
        tf = ph.text_frame
        tf.word_wrap = True
        try:
            ph.vertical_anchor = MSO_ANCHOR.MIDDLE  # type: ignore[assignment]
        except Exception:
            pass
        p = tf.paragraphs[0]
        label = spec.image_description or "Image"
        p.text = f"📷 {label}"
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            p.runs[0].font.size = _pt(14)
            p.runs[0].font.color.rgb = _hex(theme["text"])

    # Text content
    items = spec.bullets or (spec.body if isinstance(spec.body, list) else None)
    text_height = img_height
    if items:
        _add_bullets(
            slide, text_left, text_top, text_width, text_height,
            items,
            font_name=theme["body_font"],
            font_size=theme["bullet_size"],
            font_color=theme["text"],
            bullet_color=theme["accent1"],
        )
    elif spec.body and isinstance(spec.body, str):
        _add_textbox(
            slide, text_left, text_top, text_width, text_height,
            spec.body,
            font_name=theme["body_font"],
            font_size=theme["body_size"],
            font_color=theme["text"],
            line_spacing=1.4,
        )

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _build_comparison_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a comparison slide with coloured column headers."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    # Title
    if spec.title:
        _add_textbox(
            slide, left=_MARGIN_LEFT, top=_MARGIN_TOP, width=_CONTENT_WIDTH, height=0.7,
            text=spec.title,
            font_name=theme["title_font"],
            font_size=theme["title_size"],
            font_color=theme["title"],
            bold=True,
        )
        _add_accent_bar(
            slide, left=_MARGIN_LEFT, top=1.35, width=2.0, height=0.04,
            colour=theme["accent1"],
        )

    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import-untyped]

    col_width = (_CONTENT_WIDTH - 0.4) / 2
    left_x = _MARGIN_LEFT
    right_x = _MARGIN_LEFT + col_width + 0.4
    header_top = _CONTENT_TOP
    header_h = 0.55

    # Left header block (accent colour background)
    lh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _inches(left_x), _inches(header_top), _inches(col_width), _inches(header_h),
    )
    lh.fill.solid()
    lh.fill.fore_color.rgb = _hex(theme["accent1"])
    lh.line.fill.background()
    tf = lh.text_frame
    try:
        lh.vertical_anchor = MSO_ANCHOR.MIDDLE  # type: ignore[assignment]
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = spec.left_header or "Option A"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.text = p.text
    # Choose contrasting text colour for header
    run.font.color.rgb = _hex("FFFFFF")
    run.font.size = _pt(theme["heading_size"])
    run.font.bold = True
    run.font.name = theme["title_font"]

    # Right header block
    rh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _inches(right_x), _inches(header_top), _inches(col_width), _inches(header_h),
    )
    rh.fill.solid()
    rh.fill.fore_color.rgb = _hex(theme["accent2"])
    rh.line.fill.background()
    tf = rh.text_frame
    try:
        rh.vertical_anchor = MSO_ANCHOR.MIDDLE  # type: ignore[assignment]
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = spec.right_header or "Option B"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.text = p.text
    run.font.color.rgb = _hex("FFFFFF")
    run.font.size = _pt(theme["heading_size"])
    run.font.bold = True
    run.font.name = theme["title_font"]

    # Column content
    content_top = header_top + header_h + 0.25
    col_height = _SLIDE_HEIGHT_IN - content_top - _MARGIN_BOTTOM
    _render_column(slide, spec.left_column, left_x, content_top, col_width, col_height, theme)
    _render_column(slide, spec.right_column, right_x, content_top, col_width, col_height, theme)

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


def _build_blank_slide(
    prs: Any,
    spec: SlideSpec,
    theme: dict[str, Any],
    layout: Any,
) -> None:
    """Build a blank slide — only background and custom elements."""
    slide = prs.slides.add_slide(layout)
    _remove_placeholders(slide)
    bg = spec.bg_color or theme["bg"]
    _set_slide_bg(slide, bg)

    if spec.notes:
        slide.notes_slide.notes_text_frame.text = spec.notes

    if spec.elements:
        for el in spec.elements:
            _place_element(slide, el, theme)


# ---------------------------------------------------------------------------
# Image downloading
# ---------------------------------------------------------------------------


async def _download_image(url: str) -> bytes:
    """Download an image from a URL and return its bytes."""
    try:
        import httpx  # type: ignore[import-untyped]

        async with httpx.AsyncClient(follow_redirects=True, timeout=30.0) as client:
            resp = await client.get(url)
            resp.raise_for_status()
            return resp.content
    except ImportError:
        # Fall back to aiohttp if httpx is not available
        import aiohttp  # type: ignore[import-untyped]

        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=30)) as resp:
                resp.raise_for_status()
                return await resp.read()


async def _resolve_slide_images(slides: list[SlideSpec]) -> None:
    """Pre-download images from URLs and convert to base64 for all slides.

    Mutates the slide specs in-place, replacing ``image_url`` with
    ``image_b64`` so that the layout builders only need to handle base64.
    """
    for spec in slides:
        if spec.image_url and not spec.image_b64:
            try:
                img_bytes = await _download_image(spec.image_url)
                spec.image_b64 = base64.b64encode(img_bytes).decode("ascii")
                logger.info("Downloaded image for slide %r: %d bytes", spec.title, len(img_bytes))
            except Exception as exc:
                logger.warning("Failed to download image %s: %s", spec.image_url, exc)

        # Also handle image URLs in custom elements
        if spec.elements:
            for el in spec.elements:
                if el.type == "image" and not el.image_b64:
                    url = getattr(el, "image_url", None)
                    if url:
                        try:
                            img_bytes = await _download_image(url)
                            el.image_b64 = base64.b64encode(img_bytes).decode("ascii")
                        except Exception as exc:
                            logger.warning("Failed to download element image %s: %s", url, exc)


# ---------------------------------------------------------------------------
# Template loading
# ---------------------------------------------------------------------------


async def _load_template(params: CreatePptxParams) -> bytes | None:
    """Resolve the optional template from one of the supported sources."""
    if params.template_bytes_b64:
        return base64.b64decode(params.template_bytes_b64)

    if params.template_local_path:
        p = pathlib.Path(params.template_local_path).expanduser().resolve()
        if not p.is_file():
            raise FileNotFoundError(f"Template not found: {p}")
        return p.read_bytes()

    if params.template_storage_account and params.template_container and params.template_path:
        from app.connectors.blob_connector import BlobConnectorAdapter  # type: ignore[import-untyped]

        connector = BlobConnectorAdapter(
            account_name=params.template_storage_account,
            credential_type="default_azure_credential",
        )
        try:
            return await connector.download_blob(params.template_container, params.template_path)
        finally:
            await connector.close()

    return None


# ---------------------------------------------------------------------------
# Output writing
# ---------------------------------------------------------------------------


async def _save_output(
    pptx_bytes: bytes,
    params: CreatePptxParams,
) -> dict[str, str | None]:
    """Write the PPTX to local path and/or blob storage. Returns paths."""
    result: dict[str, str | None] = {"local_path": None, "blob_path": None}

    if params.output_local_path:
        p = pathlib.Path(params.output_local_path).expanduser().resolve()
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_bytes(pptx_bytes)
        result["local_path"] = str(p)

    if params.output_storage_account and params.output_container and params.output_path:
        from app.connectors.blob_connector import BlobConnectorAdapter  # type: ignore[import-untyped]

        connector = BlobConnectorAdapter(
            account_name=params.output_storage_account,
            credential_type="default_azure_credential",
        )
        try:
            await connector.upload_blob(
                params.output_container,
                params.output_path,
                pptx_bytes,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            result["blob_path"] = (
                f"https://{params.output_storage_account}.blob.core.windows.net"
                f"/{params.output_container}/{params.output_path}"
            )
        finally:
            await connector.close()

    return result


# ---------------------------------------------------------------------------
# Layout selection & placeholder cleanup
# ---------------------------------------------------------------------------


def _find_blank_layout(prs: Any) -> Any:
    """Return the slide layout with the fewest/smallest placeholders.

    Searches by name first (common names for blank layouts across templates),
    then falls back to whichever layout has no content placeholders (idx 0/1).
    """
    # Try well-known blank layout names
    for layout in prs.slide_layouts:
        if layout.name.lower().strip() in ("blank", "leer", "vide", "en blanco"):
            return layout

    # Fall back: pick the layout with the fewest placeholders whose idx < 10
    # (idx 10-12 are date/footer/slide-number — harmless)
    best = None
    best_count = 999
    for layout in prs.slide_layouts:
        content_phs = [
            ph for ph in layout.placeholders
            if ph.placeholder_format.idx < 10
        ]
        if len(content_phs) < best_count:
            best_count = len(content_phs)
            best = layout
    return best or prs.slide_layouts[-1]


def _remove_placeholders(slide: Any) -> None:
    """Remove all inherited placeholder shapes from a slide.

    Layout placeholders (title, body, etc.) are added automatically when a
    slide is created from a layout.  Since we draw our own styled text boxes,
    these empty placeholders overlap and hide our content.  This helper
    deletes them from the slide's XML shape tree.
    """
    from pptx.oxml.ns import qn  # type: ignore[import-untyped]

    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    for sp in list(spTree.iterchildren(qn("p:sp"))):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        # A placeholder has a <p:ph> child element
        if nvPr.find(qn("p:ph")) is not None:
            spTree.remove(sp)


# ---------------------------------------------------------------------------
# Main tool
# ---------------------------------------------------------------------------

_LAYOUT_BUILDERS = {
    "title": _build_title_slide,
    "section": _build_section_slide,
    "content": _build_content_slide,
    "two_column": _build_two_column_slide,
    "comparison": _build_comparison_slide,
    "blank": _build_blank_slide,
}


@define_tool(
    name="create_pptx",
    description=(
        "Create or edit a PowerPoint (PPTX) presentation with professional styling. "
        "Supports multiple slide layouts (title, content, section divider, two-column, "
        "image+text, comparison, blank), built-in colour themes (corporate, modern_dark, "
        "creative, minimal), custom elements (text boxes, images, shapes with content), "
        "and optional template files as a starting point. Returns the presentation as "
        "base64 and optionally saves to local disk or Azure Blob Storage."
    ),
    parameters_model=CreatePptxParams,
)
async def create_pptx(params: CreatePptxParams, context: dict | None = None) -> dict[str, Any]:
    """Build a PPTX presentation from the given slide specifications."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches  # type: ignore[import-untyped]

        # 1. Load template or create blank presentation
        template_bytes = await _load_template(params)
        if template_bytes:
            prs = Presentation(io.BytesIO(template_bytes))
        else:
            prs = Presentation()

        # 2. Set slide dimensions
        prs.slide_width = Inches(params.slide_width)
        prs.slide_height = Inches(params.slide_height)

        # 3. Resolve theme
        theme = _resolve_theme(params)

        # 4. Find a truly blank layout (no content placeholders)
        blank_layout = _find_blank_layout(prs)

        # 5. Pre-download any images referenced by URL
        await _resolve_slide_images(params.slides)

        # 6. Build each slide
        for slide_spec in params.slides:
            layout_name = slide_spec.layout

            if layout_name in ("image_left", "image_right"):
                side: Literal["left", "right"] = "left" if layout_name == "image_left" else "right"
                _build_image_slide(prs, slide_spec, theme, blank_layout, image_side=side)
            elif layout_name in _LAYOUT_BUILDERS:
                _LAYOUT_BUILDERS[layout_name](prs, slide_spec, theme, blank_layout)
            else:
                _build_content_slide(prs, slide_spec, theme, blank_layout)

        # 7. Serialise to bytes
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        # 8. Save to destinations
        destinations = await _save_output(pptx_bytes, params)

        # 9. Build response
        filename = params.output_filename or "presentation.pptx"
        return {
            "filename": filename,
            "slide_count": len(params.slides),
            "theme": params.theme,
            "file_b64": base64.b64encode(pptx_bytes).decode("ascii"),
            "file_size_bytes": len(pptx_bytes),
            "output_local_path": destinations.get("local_path"),
            "output_blob_url": destinations.get("blob_path"),
        }

    except ImportError as exc:
        logger.exception("create_pptx failed — missing dependency")
        return {
            "error": f"Missing required dependency: {exc}. Ensure 'python-pptx' is installed (pip install python-pptx).",
            "filename": params.output_filename or "presentation.pptx",
        }
    except FileNotFoundError as exc:
        logger.exception("create_pptx failed — template not found")
        return {
            "error": f"Template file not found: {exc}",
            "filename": params.output_filename or "presentation.pptx",
        }
    except Exception as exc:
        logger.exception("create_pptx failed")
        return {
            "error": (
                f"Failed to create presentation: {type(exc).__name__}: {exc}. "
                "Check that slide specifications are valid (layouts, content blocks, "
                "image data, and element positions) and that any template files are "
                "accessible."
            ),
            "filename": params.output_filename or "presentation.pptx",
        }
