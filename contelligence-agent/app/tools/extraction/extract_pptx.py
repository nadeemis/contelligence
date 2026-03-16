"""Tool for extracting slide content from PowerPoint PPTX files."""

from __future__ import annotations

import base64
import io
import logging
import pathlib
from typing import Any, Literal

from pydantic import BaseModel, Field

from app.core.tool_registry import define_tool

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Parameters
# ---------------------------------------------------------------------------

class ExtractPptxParams(BaseModel):
    """Parameters for the extract_pptx tool."""

    file_bytes: bytes | None = Field(
        None,
        description="Raw bytes of the PPTX file. Preferred when the caller already has the content in memory.",
    )
    file_bytes_b64: str | None = Field(
        None,
        description=(
            "Base64-encoded PPTX file content. When provided, the file is "
            "parsed directly from the supplied bytes instead of reading from "
            "Azure Blob Storage."
        ),
    )
    local_path: str | None = Field(
        None,
        description="Absolute or relative path to a PPTX file on the local filesystem.",
    )
    storage_account: str | None = Field(
        default=None,
        description=(
            "Azure Storage account name to connect to. "
            "When omitted the default account configured for the agent is used. "
            "Provide this to write to a different storage account "
            "(authentication uses DefaultAzureCredential)."
        ),
    )
    container: str | None = Field(
        None, description="Azure Blob Storage container name (used when file_bytes_b64 is not provided)."
    )
    path: str | None = Field(
        None, description="Blob path to the PPTX file (used when file_bytes_b64 is not provided)."
    )
    format: Literal["markdown", "json"] = Field(
        "markdown",
        description="Output format: 'markdown' returns a rendered markdown document, 'json' returns structured slide data.",
    )
    filename: str | None = Field(
        None,
        description="Optional descriptive filename for the result metadata.",
    )


# ---------------------------------------------------------------------------
# Tool definition
# ---------------------------------------------------------------------------

@define_tool(
    name="extract_pptx",
    description=(
        "Extract slide content from a PowerPoint PPTX file. Accepts the file "
        "as raw bytes (file_bytes), base64-encoded bytes (file_bytes_b64), a "
        "local filesystem path (local_path), or reads from Azure Blob Storage "
        "(container + path). Returns a markdown document (default) or "
        "structured JSON. Uses python-pptx to read titles, text, notes, "
        "and shape types from each slide."
    ),
    parameters_model=ExtractPptxParams,
)
async def extract_pptx(params: ExtractPptxParams, context: dict) -> dict[str, Any]:
    """Extract slide content from a PPTX provided as bytes or fetched from blob storage."""
    try:
        from pptx import Presentation

        # Resolve the file bytes from the supplied source.
        if params.file_bytes:
            logger.info("Parsing PPTX from supplied raw bytes")
            pptx_bytes: bytes = params.file_bytes
        elif params.file_bytes_b64:
            logger.info("Parsing PPTX from base64-encoded bytes")
            pptx_bytes = base64.b64decode(params.file_bytes_b64)
        elif params.local_path:
            resolved = pathlib.Path(params.local_path).expanduser().resolve()
            if not resolved.is_file():
                return {
                    "error": f"Local file not found: {resolved}",
                    "filename": params.filename or params.local_path,
                }
            logger.info("Reading PPTX from local path %s", resolved)
            pptx_bytes = resolved.read_bytes()
        elif params.storage_account and params.container and params.path:
            from app.connectors.blob_connector import BlobConnectorAdapter
            try:
                connector = BlobConnectorAdapter(
                    account_name=params.storage_account,
                    credential_type="default_azure_credential",
                )
                logger.info("Downloading PPTX blob %s/%s", params.container, params.path)
                pptx_bytes = await connector.download_blob(params.container, params.path)
            except Exception as blob_err:
                return {
                    "error": f"Failed to download PPTX from blob storage: {blob_err}",
                    "filename": params.filename or params.path,
                }
            finally:
                if connector is not None:
                    await connector.close()
        else:
            return {
                "error": (
                    "Provide one of: file_bytes, file_bytes_b64, local_path, "
                    "or storage_account + container + path."
                ),
                "filename": params.filename or "",
            }

        prs = Presentation(io.BytesIO(pptx_bytes))

        slides_data: list[dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            # --- title ---
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text

            # --- text from all shapes ---
            texts: list[str] = []
            shape_types: list[str] = []
            for shape in slide.shapes:
                shape_types.append(str(shape.shape_type))
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            texts.append(para_text)

            # --- notes ---
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text

            slides_data.append(
                {
                    "number": idx,
                    "title": title,
                    "text": "\n".join(texts),
                    "notes": notes,
                    "shapes": shape_types,
                }
            )

        label = params.filename or params.local_path or params.path or "unknown.pptx"

        if params.format == "markdown":
            return {
                "filename": label,
                "slide_count": len(slides_data),
                "content": _slides_to_markdown(label, slides_data),
            }

        return {
            "filename": label,
            "slide_count": len(slides_data),
            "slides": slides_data,
        }

    except Exception as exc:
        label = params.filename or params.local_path or params.path or "unknown.pptx"
        logger.exception("extract_pptx failed for %s", label)
        return {"error": str(exc), "filename": label}


# ---------------------------------------------------------------------------
# Markdown renderer
# ---------------------------------------------------------------------------

def _slides_to_markdown(label: str, slides: list[dict[str, Any]]) -> str:
    """Convert structured slide data into a readable markdown document."""
    parts: list[str] = [f"# {label}", ""]

    for slide in slides:
        num = slide["number"]
        title = slide.get("title", "")
        text = slide.get("text", "")
        notes = slide.get("notes", "")

        heading = f"## Slide {num}"
        if title:
            heading += f" — {title}"
        parts.append(heading)
        parts.append("")

        if text:
            parts.append(text)
            parts.append("")

        if notes:
            parts.append("> **Notes:** " + notes.replace("\n", " "))
            parts.append("")

    return "\n".join(parts).rstrip() + "\n"
